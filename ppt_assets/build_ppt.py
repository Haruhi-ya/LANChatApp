# -*- coding: utf-8 -*-
"""生成汇报答辩PPT：局域网聊天室（即时通讯工具）"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ---------- 常量 ----------
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
FONT = "微软雅黑"
NAVY   = RGBColor(0x1F, 0x4E, 0x79)   # 主色 深蓝
BLUE   = RGBColor(0x2E, 0x75, 0xB6)   # 辅色
SKY    = RGBColor(0x5B, 0x9B, 0xD5)
ORANGE = RGBColor(0xE8, 0xA3, 0x3D)   # 强调
LIGHT  = RGBColor(0xF2, 0xF6, 0xFB)   # 浅底
GRAY   = RGBColor(0x59, 0x5F, 0x69)
DARK   = RGBColor(0x26, 0x2B, 0x33)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x2E, 0x8B, 0x57)

prs = Presentation()
prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
BLANK = prs.slide_layouts[6]
SHOT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "shots")
ARCH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "架构图-新.png")

# ---------- 基础工具 ----------

def add_slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp

def txt(slide, x, y, w, h, text, size=14, color=DARK, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        f = r.font
        f.name = FONT
        f.size = Pt(size)
        f.bold = bold
        f.color.rgb = color
        # 中文字体兜底
        rPr = r._r.get_or_add_rPr()
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = rPr.makeelement(qn('a:ea'), {})
            rPr.append(ea)
        ea.set('typeface', FONT)
    return tb

def bullet_box(slide, x, y, w, h, items, size=13.5, gap=6):
    """items: [(文本, 层级, 加粗)] 或纯字符串列表"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        if isinstance(it, tuple):
            text, lvl, bold = it
        else:
            text, lvl, bold = it, 0, False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        if lvl > 0:
            r0 = p.add_run(); r0.text = "  " * lvl + "▪ "
            r0.font.size = Pt(size); r0.font.color.rgb = BLUE; r0.font.bold = True; r0.font.name = FONT
        else:
            r0 = p.add_run(); r0.text = "● "
            r0.font.size = Pt(size); r0.font.color.rgb = ORANGE; r0.font.bold = True; r0.font.name = FONT
        r = p.add_run(); r.text = text
        f = r.font; f.name = FONT; f.size = Pt(size); f.bold = bold; f.color.rgb = DARK
    return tb

def page_header(slide, title, subtitle=""):
    """内页顶部标题条"""
    rect(slide, 0, 0, SLIDE_W, Inches(0.92), NAVY)
    rect(slide, 0, Inches(0.92), SLIDE_W, Inches(0.045), ORANGE)
    txt(slide, Inches(0.55), Inches(0.10), Inches(9.5), Inches(0.75), title,
        size=25, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        txt(slide, Inches(9.4), Inches(0.10), Inches(3.5), Inches(0.75), subtitle,
            size=12, color=RGBColor(0xBD, 0xD7, 0xEE), align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    # 页码
    n = len(prs.slides.__iter__.__self__._sldIdLst) if False else None

def footer(slide, page_no):
    rect(slide, 0, Inches(7.18), SLIDE_W, Inches(0.32), LIGHT)
    txt(slide, Inches(12.3), Inches(7.19), Inches(0.8), Inches(0.3),
        str(page_no), size=10, color=GRAY, align=PP_ALIGN.RIGHT)

def card(slide, x, y, w, h, title, body, title_color=NAVY, body_size=12.5):
    rect(slide, x, y, w, h, LIGHT, line=RGBColor(0xD5, 0xE2, 0xF0), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(slide, x + Inches(0.14), y + Inches(0.07), w - Inches(0.28), Inches(0.35),
        title, size=14, color=title_color, bold=True)
    txt(slide, x + Inches(0.14), y + Inches(0.42), w - Inches(0.28), h - Inches(0.5),
        body, size=body_size, color=DARK, line_spacing=1.12)

# ---------- 1. 封面 ----------
s = add_slide()
rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
rect(s, 0, Inches(2.35), SLIDE_W, Inches(0.03), ORANGE)
txt(s, 0, Inches(1.05), SLIDE_W, Inches(1.2), "局域网聊天室（即时通讯工具）",
    size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, 0, Inches(2.55), SLIDE_W, Inches(0.55), "软件工程课程设计 · 汇报答辩",
    size=17, color=RGBColor(0xBD, 0xD7, 0xEE), align=PP_ALIGN.CENTER)

# 成员信息表（用色块+文本模拟表格）
mem = [
    ("赵翔", "25112002136", "网络通信与服务端"),
    ("胡鸣", "25112002139", "数据库设计与数据层"),
    ("茹桂堂", "25112002124", "客户端界面与交互"),
]
x0, y0 = Inches(2.2), Inches(3.55)
for i, (nm, sid, duty) in enumerate(mem):
    bx = x0 + Inches(i * 3.0)
    rect(s, bx, y0, Inches(2.8), Inches(1.15), RGBColor(0x2A, 0x63, 0x97),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, bx, y0 + Inches(0.12), Inches(2.8), Inches(0.4), nm + "  " + sid,
        size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, bx, y0 + Inches(0.55), Inches(2.8), Inches(0.5), duty,
        size=12, color=RGBColor(0xBD, 0xD7, 0xEE), align=PP_ALIGN.CENTER)
txt(s, 0, Inches(5.05), SLIDE_W, Inches(0.4), "分工：网络通信 / 界面交互 / 数据存储 三条技术线并行协作",
    size=13, color=RGBColor(0x9D, 0xC3, 0xE6), align=PP_ALIGN.CENTER)
txt(s, 0, Inches(5.55), SLIDE_W, Inches(0.4), "2026 年 8 月 · 6号楼404",
    size=15, color=WHITE, align=PP_ALIGN.CENTER)

# ---------- 2. 目录 ----------
s = add_slide()
page_header(s, "目 录", "CONTENTS")
dirs = [
    ("01", "项目背景", "课题来源与项目目标"),
    ("02", "软件需求", "六大功能域 · 非功能需求"),
    ("03", "概要设计", "系统架构 · 模块划分 · 数据库"),
    ("04", "详细设计", "通信协议 · 关键技术 · 安全"),
    ("05", "系统实现", "界面展示与核心功能"),
    ("06", "系统测试", "测试环境 · 用例与结论"),
]
for i, (no, t, d) in enumerate(dirs):
    col, row = i % 2, i // 2
    bx = Inches(0.9 + col * 6.1)
    by = Inches(1.5 + row * 1.85)
    rect(s, bx, by, Inches(5.6), Inches(1.5), LIGHT, line=RGBColor(0xD5, 0xE2, 0xF0),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, bx + Inches(0.25), by + Inches(0.28), Inches(1.0), Inches(0.9), no,
        size=34, color=ORANGE, bold=True)
    txt(s, bx + Inches(1.25), by + Inches(0.22), Inches(4.2), Inches(0.5), t,
        size=20, color=NAVY, bold=True)
    txt(s, bx + Inches(1.25), by + Inches(0.78), Inches(4.2), Inches(0.5), d,
        size=12.5, color=GRAY)
footer(s, 2)

# ---------- 3. 项目背景 ----------
s = add_slide()
page_header(s, "项目背景")
card(s, Inches(0.6), Inches(1.3), Inches(6.0), Inches(2.6), "课题来源",
     "局域网聊天室是软件工程课程设计的典型课题：在网络编程基础上，综合运用 Java 多线程、Socket 通信、Swing 界面与数据库技术，在局域网内实现无外网依赖的即时通讯。\n\n本组选择该课题，旨在完整实践「网络层、界面层、数据层」三层分离的团队协作开发。")
card(s, Inches(6.9), Inches(1.3), Inches(5.9), Inches(2.6), "项目目标",
     "基于 Java 语言设计实现一套完整的局域网聊天室系统：\n\n• 客户端：Swing 图形界面框架\n• 服务端：Socket 多线程模型（TCP 8080）\n• 数据：MySQL 持久化，JDBC 访问\n• 功能：注册登录、公共聊天室、一对一私聊、历史记录、消息撤回、记录搜索、管理员管理、自定义头像与图片消息")
card(s, Inches(0.6), Inches(4.15), Inches(6.0), Inches(2.35), "技术要点",
     "• 一连接一线程的多线程服务端\n• 自定义文本行通信协议\n• JTable 自绘气泡消息列表\n• Base64 图片消息传输与校验\n• 双份行模型的私聊记录存储")
card(s, Inches(6.9), Inches(4.15), Inches(5.9), Inches(2.35), "开发方式",
     "• 三人按「网络通信、界面交互、数据存储」三大技术线分工\n• 接口先行：开工前统一协议格式与回调接口\n• 联合调试：双客户端、图片与头像功能联测\n• 统一使用 IDEA + Git 版本管理")
footer(s, 3)

# ---------- 4. 软件需求（功能） ----------
s = add_slide()
page_header(s, "软件需求 —— 六大功能域")
feats = [
    ("用户管理", "注册登录、角色权限\n（admin / user）"),
    ("公共聊天室", "消息实时显示、历史\n回放、图片消息"),
    ("私聊", "一对一私聊、离线投递\n未读提示、历史搜索"),
    ("消息管理", "两分钟撤回、记录清空\n关键字搜索"),
    ("管理员功能", "踢出用户、封禁账号\n清空公共记录"),
    ("个性化功能", "自定义头像、图片消息\n自动压缩适配带宽"),
]
for i, (t, d) in enumerate(feats):
    col, row = i % 3, i // 3
    bx = Inches(0.6 + col * 4.15)
    by = Inches(1.35 + row * 2.5)
    rect(s, bx, by, Inches(3.85), Inches(2.15), LIGHT, line=RGBColor(0xD5, 0xE2, 0xF0),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, bx, by, Inches(0.09), Inches(2.15), ORANGE)
    txt(s, bx + Inches(0.3), by + Inches(0.18), Inches(3.4), Inches(0.45), t,
        size=16.5, color=NAVY, bold=True)
    txt(s, bx + Inches(0.3), by + Inches(0.75), Inches(3.4), Inches(1.3), d,
        size=12.5, color=DARK, line_spacing=1.25)
card(s, Inches(0.6), Inches(6.15), Inches(12.2), Inches(0.82), "", "", body_size=12)
txt(s, Inches(0.85), Inches(6.28), Inches(11.8), Inches(0.6),
    "非功能需求：连接失败、服务器断开、历史加载超时等异常场景均有明确提示与兜底处理，客户端不出现无响应状态。",
    size=12.5, color=DARK)
footer(s, 4)

# ---------- 5. 概要设计：系统总体架构（原生形状绘制，可编辑） ----------
s = add_slide()
page_header(s, "概要设计 —— 系统总体架构")

def arrow_line(slide, x1, y1, x2, y2, color=BLUE, width=1.6, both=False):
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    attrs = {'type': 'arrow', 'w': 'med', 'len': 'med'}
    if both:
        ln.append(ln.makeelement(qn('a:headEnd'), attrs))
    ln.append(ln.makeelement(qn('a:tailEnd'), attrs))
    return conn

def mini_box(slide, x, y, w, h, name, sub, fc=LIGHT, ec=SKY, name_size=10.5, sub_size=8, name_color=NAVY):
    rect(slide, x, y, w, h, fc, line=ec, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    if sub:
        txt(slide, x + Inches(0.05), y + Inches(0.08), w - Inches(0.1), Inches(0.3),
            name, size=name_size, color=name_color, bold=True, align=PP_ALIGN.CENTER)
        txt(slide, x + Inches(0.05), y + Inches(0.38), w - Inches(0.1), Inches(0.3),
            sub, size=sub_size, color=GRAY, align=PP_ALIGN.CENTER)
    else:
        txt(slide, x, y, w, h, name, size=name_size, color=name_color, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def container(slide, x, y, w, h, title):
    rect(slide, x, y, w, h, LIGHT, line=SKY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(slide, x + Inches(0.01), y + Inches(0.01), w - Inches(0.02), Inches(0.34), NAVY,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(slide, x, y + Inches(0.02), w, Inches(0.3), title, size=12.5, color=WHITE,
        bold=True, align=PP_ALIGN.CENTER)

CX, CY = Inches(0.45), Inches(1.42)
CW = Inches(4.85)
SX, SY = Inches(8.0), Inches(1.42)
SWd = Inches(4.85)
CH = Inches(4.15)

# ---- 客户端容器 ----
container(s, CX, CY, CW, CH, "客 户 端（Swing 界面 + 网络层）")
txt(s, CX + Inches(0.4), CY + Inches(0.46), Inches(4.0), Inches(0.2),
    "界面层（事件分发线程 EDT 更新）", size=9, color=BLUE, bold=True)
mods = [("登录 / 注册", "chatEntryUI"), ("主聊天窗口", "公共聊天·好友·管理"),
        ("私聊 / 对话框", "私聊·好友申请")]
for i, (nm, sub) in enumerate(mods):
    bx = CX + Inches(0.4 + i * 1.48)
    mini_box(s, bx, CY + Inches(0.68), Inches(1.38), Inches(0.78), nm, sub)
    # 界面层 → 网络层 箭头
    arrow_line(s, bx + Inches(0.69), CY + Inches(1.46), bx + Inches(0.69), CY + Inches(1.70), width=1.1)
rect(s, CX + Inches(0.4), CY + Inches(1.72), Inches(4.05), Inches(0.6), BLUE,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, CX + Inches(0.4), CY + Inches(1.79), Inches(4.05), Inches(0.45), "网络层 chatClient",
    size=10.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, CX + Inches(0.4), CY + Inches(2.48), Inches(4.05), Inches(0.85),
    "独立接收线程 · Listener 回调\nSwingUtilities.invokeLater 切回 EDT",
    size=8.5, color=GRAY, align=PP_ALIGN.CENTER, line_spacing=1.15)
txt(s, CX + Inches(0.4), CY + Inches(3.5), Inches(4.05), Inches(0.5),
    "发送方法 sendLine 统一出口\n过滤 \\r\\n，防止换行注入拆命令",
    size=8.5, color=GRAY, align=PP_ALIGN.CENTER, line_spacing=1.15)

# ---- 服务端容器 ----
container(s, SX, SY, SWd, CH, "服 务 端（Socket 多线程）")
mini_box(s, SX + Inches(0.4), SY + Inches(0.5), Inches(1.95), Inches(0.85),
         "连接管理", "主线程 accept\n一连接一线程")
mini_box(s, SX + Inches(2.5), SY + Inches(0.5), Inches(1.95), Inches(0.85),
         "协议分发", "handleMessage\n29 类命令")
mini_box(s, SX + Inches(0.4), SY + Inches(1.5), Inches(4.05), Inches(0.85),
         "消息管道", "清洗 → 魔数校验 → 入库 → 广播 → @提醒")
txt(s, SX + Inches(0.4), SY + Inches(2.35), Inches(4.05), Inches(0.25),
    "先入库后广播，历史回放与实时广播一致", size=8, color=GRAY, align=PP_ALIGN.CENTER)
mini_box(s, SX + Inches(0.4), SY + Inches(2.62), Inches(1.95), Inches(0.85),
         "会话与权限", "登录校验\n踢人 · 封禁")
mini_box(s, SX + Inches(2.5), SY + Inches(2.62), Inches(1.95), Inches(0.85),
         "好友管理", "好友申请\n列表 · 删除")
txt(s, SX + Inches(0.4), SY + Inches(3.58), Inches(4.05), Inches(0.45),
    "登录校验 → 权限校验（isAdmin）\n未登录连接忽略一切命令",
    size=8.5, color=GRAY, align=PP_ALIGN.CENTER, line_spacing=1.15)

# ---- TCP 双向箭头 ----
arrow_line(s, Inches(5.38), Inches(3.05), Inches(7.92), Inches(3.05),
           color=ORANGE, width=2.0, both=True)
txt(s, Inches(5.38), Inches(1.85), Inches(2.54), Inches(0.6), "TCP Socket\n自定义文本行协议",
    size=10, color=NAVY, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.05)
txt(s, Inches(5.38), Inches(3.25), Inches(2.54), Inches(0.55), "UTF-8 · 冒号分隔\n29 类命令",
    size=8.5, color=GRAY, align=PP_ALIGN.CENTER, line_spacing=1.1)

# ---- JDBC 箭头 ----
arrow_line(s, Inches(10.42), Inches(5.57), Inches(10.42), Inches(5.93), color=BLUE, width=1.8)
txt(s, Inches(10.02), Inches(5.63), Inches(0.8), Inches(0.28), "JDBC", size=9, color=NAVY,
    bold=True, align=PP_ALIGN.CENTER)

# ---- 数据库容器 ----
rect(s, Inches(0.45), Inches(5.95), Inches(12.4), Inches(0.95), LIGHT,
     line=SKY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, Inches(0.45), Inches(5.95), Inches(1.5), Inches(0.95), NAVY)
txt(s, Inches(0.45), Inches(6.12), Inches(1.5), Inches(0.6), "MySQL\n数据库 (lanchat)",
    size=9.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.0)
dbs = [("Users", "用户·角色·头像"), ("PublicMessages", "公共消息·撤回ID"),
       ("PrivateMessages", "私聊双份行·未读"), ("Friendships", "好友关系"),
       ("FriendRequests", "好友申请")]
for i, (nm, sub) in enumerate(dbs):
    bx = Inches(2.1 + i * 2.06)
    mini_box(s, bx, Inches(6.03), Inches(1.96), Inches(0.78), nm, sub,
             name_size=10, sub_size=8)
txt(s, Inches(0.45), Inches(6.98), Inches(12.4), Inches(0.2),
    "启动时自动建库建表与结构迁移（补 role / avatar 列 · TEXT→MEDIUMTEXT · client_id 索引）",
    size=9, color=GRAY, align=PP_ALIGN.CENTER)
footer(s, 5)

# ---------- 6. 概要设计：模块与数据库 ----------
s = add_slide()
page_header(s, "概要设计 —— 模块划分与数据库设计")
card(s, Inches(0.6), Inches(1.22), Inches(6.0), Inches(2.28), "客户端模块",
     "• 网络层 chatClient：独立接收线程 + 监听器回调\n• 登录界面 chatEntryUI：注册/登录/服务器配置\n• 主聊天窗口 clientChatUI：公共聊天、好友列表、未读徽标、头像入口\n• 私聊窗口 privateChatUI：独立会话、搜索、清空\n• 气泡消息列表 bubbleChatList：单列 JTable 自绘")
card(s, Inches(6.9), Inches(1.22), Inches(5.9), Inches(2.28), "服务端模块",
     "• 连接管理：主线程 accept，一连接一线程\n• 协议分发：按命令前缀分发，先登录校验后执行\n• 消息管道：清洗 → 入库 → 广播 → @检测\n• 会话与权限：登录校验、踢人、封禁；好友管理\n• 数据访问 dbManager：建库建表、结构自动迁移、增删改查")
# 数据库表（5 张，含好友功能两张新表）
tbl = s.shapes.add_table(6, 5, Inches(0.6), Inches(3.68), Inches(12.2), Inches(2.85)).table
tbl.columns[0].width = Inches(2.5)
tbl.columns[1].width = Inches(3.3)
tbl.columns[2].width = Inches(2.7)
tbl.columns[3].width = Inches(1.5)
tbl.columns[4].width = Inches(2.2)
headers = ["数据表", "关键字段", "用途", "要点"]
rows = [
    ("Users 用户表", "username(主键) / password / role / avatar", "用户账号、角色、头像", "头像 MEDIUMBLOB"),
    ("PublicMessages 公共消息表", "id / send_time / sender / content / client_id", "公共聊天历史", "client_id 供撤回定位"),
    ("PrivateMessages 私聊消息表", "owner / peer / sender / is_read / content / client_id", "私聊历史与未读", "双份行模型，各自清空互不影响"),
    ("Friendships 好友关系表", "user_a / user_b / created_at", "双向好友关系", "user_a<user_b 规范化一行"),
    ("FriendRequests 好友申请表", "from_user / to_user / created_at", "待处理好友申请", "UNIQUE 防重复申请"),
]
# 表头
for j, h in enumerate(headers):
    c = tbl.cell(0, j)
    c.fill.solid(); c.fill.fore_color.rgb = NAVY
    c.text_frame.text = h
    for p in c.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(11.5); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
for i, row in enumerate(rows):
    for j, v in enumerate(row):
        c = tbl.cell(i + 1, j)
        c.fill.solid()
        c.fill.fore_color.rgb = LIGHT if i % 2 == 0 else WHITE
        c.text_frame.text = v
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(10.5); r.font.color.rgb = DARK; r.font.name = FONT
txt(s, Inches(0.6), Inches(6.62), Inches(12.2), Inches(0.45),
    "系统启动时自动建库建表并完成结构迁移（补 role / avatar 列、TEXT→MEDIUMTEXT 扩容、加 client_id 索引）。",
    size=11, color=GRAY)
footer(s, 6)

# ---------- 7. 详细设计：协议 ----------
s = add_slide()
page_header(s, "详细设计 —— 通信协议")
card(s, Inches(0.6), Inches(1.25), Inches(6.1), Inches(2.85), "协议设计要点",
     "• 自定义文本行协议：一行一条命令，UTF-8 编码，冒号分隔字段\n• 内容字段固定放最后：indexOf(':') 只切固定字段，内容可含冒号\n• 保留字符约束：冒号、逗号禁止出现在用户名中\n• 时间戳用 epoch 毫秒纯数字，规避 HH:mm:ss 冒号问题\n• 消息 ID 由服务端生成 UUID，不接受客户端传入，防伪造撤回")
card(s, Inches(6.95), Inches(1.25), Inches(5.85), Inches(2.85), "核心命令（节选）",
     "客户端→服务端：\n• REGISTER / LOGIN —— 注册登录\n• MSG / PM —— 公共 / 私聊消息\n• RECALL / PMRECALL —— 撤回\n• KICK / BAN / CLEARPUBLIC —— 管理\n• HIST / UNREAD / SEARCHPUB —— 历史未读搜索\n• GETAVATAR / SETAVATAR —— 头像\n服务端→客户端：\n• LOGINOK / MSG / PMMSG / SYSTEM\n• USERS / HISTBEGIN·ITEM·END / RECALLED 等 40 余类")
card(s, Inches(0.6), Inches(4.35), Inches(6.1), Inches(2.4), "消息管道一致性",
     "内容型消息依次经过：\n长度与格式清洗 → 图片合法性校验 → 数据库保存 → 实时广播 → @提醒扫描\n\n先入库后广播：保证历史回放与实时广播的顺序完全一致，入库失败则不广播。")
card(s, Inches(6.95), Inches(4.35), Inches(5.85), Inches(2.4), "图片消息传输",
     "• 图片压缩至长边 ≤1280px、体积 ≤1MB\n• 以 [IMG]+Base64 内容自描述格式传输与存储，协议命令零扩展\n• 服务端校验 Base64 合法性与图片魔数（JPEG/PNG 文件头），非法数据直接丢弃\n• 内容列 MEDIUMTEXT 容纳 Base64 体积")
footer(s, 7)

# ---------- 8. 详细设计：关键技术 ----------
s = add_slide()
page_header(s, "详细设计 —— 关键技术")
techs = [
    ("JTable 可变行高气泡渲染", "单列 JTable 自绘：prepareRenderer 按内容计算行高，LineBreakMeasurer 按字符边界断行，中文/英文/表情换行正确、不拆散代理对；绘制在渲染器内完成，避免事件派发问题。"),
    ("头像按需拉取与版本控制", "头像不随列表广播，绘制时按需拉取，内存缓存 + 负缓存防重复请求；变更时广播通知清除缓存；版本号机制丢弃陈旧响应，快速连换头像最终一致。"),
    ("消息撤回机制", "服务端生成全局唯一消息 ID；撤回校验 2 分钟时间窗口 + 操作者权限（普通用户仅自己的消息，管理员任意）；通过后广播撤回事件，全端同步替换为撤回提示。"),
    ("多线程与并发控制", "服务端一连接一线程，在线表 ConcurrentHashMap；客户端独立接收线程 + SwingUtilities.invokeLater 切回 EDT；历史加载期间实时消息先缓冲后按序补渲染，避免重复乱序。"),
]
for i, (t, d) in enumerate(techs):
    col, row = i % 2, i // 2
    bx = Inches(0.6 + col * 6.25)
    by = Inches(1.3 + row * 2.75)
    rect(s, bx, by, Inches(5.95), Inches(2.45), LIGHT, line=RGBColor(0xD5, 0xE2, 0xF0),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, bx, by, Inches(0.09), Inches(2.45), BLUE)
    txt(s, bx + Inches(0.3), by + Inches(0.15), Inches(5.5), Inches(0.4), t,
        size=15, color=NAVY, bold=True)
    txt(s, bx + Inches(0.3), by + Inches(0.62), Inches(5.45), Inches(1.75), d,
        size=12, color=DARK, line_spacing=1.18)
footer(s, 8)

# ---------- 9. 系统实现：界面展示1 ----------
s = add_slide()
page_header(s, "系统实现 —— 界面展示（一）")
pic = s.shapes.add_picture(os.path.join(SHOT, "01-登录界面.png"), Inches(0.5), Inches(1.35), height=Inches(4.6))
pic = s.shapes.add_picture(os.path.join(SHOT, "02-主聊天窗口-公共聊天室.png"), Inches(4.45), Inches(1.55), width=Inches(6.0))
txt(s, Inches(1.25), Inches(6.15), Inches(2.2), Inches(0.4), "登录 / 注册界面",
    size=11.5, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, Inches(6.0), Inches(6.15), Inches(3.6), Inches(0.4), "主聊天窗口（公共聊天室）",
    size=11.5, color=GRAY, align=PP_ALIGN.CENTER)
card(s, Inches(10.65), Inches(1.3), Inches(2.2), Inches(5.3), "功能入口",
     "• 在线用户列表（未读红点）\n• 双击发起私聊\n• 好友申请与管理\n• @提醒、图片发送\n• 自定义头像\n• 消息撤回\n• 管理员入口", body_size=10.5)
footer(s, 9)

# ---------- 10. 系统实现：界面展示2 ----------
s = add_slide()
page_header(s, "系统实现 —— 界面展示（二）")
pic = s.shapes.add_picture(os.path.join(SHOT, "04-私聊窗口.png"), Inches(0.5), Inches(1.35), height=Inches(4.6))
txt(s, Inches(1.1), Inches(6.1), Inches(2.6), Inches(0.4), "一对一私聊窗口",
    size=11.5, color=GRAY, align=PP_ALIGN.CENTER)
pic = s.shapes.add_picture(os.path.join(SHOT, "05-管理员-查看私聊记录.png"), Inches(4.75), Inches(1.35), height=Inches(4.6))
txt(s, Inches(5.15), Inches(6.1), Inches(3.0), Inches(0.4), "管理员查看任意用户私聊记录",
    size=11.5, color=GRAY, align=PP_ALIGN.CENTER)
card(s, Inches(9.0), Inches(1.3), Inches(3.85), Inches(2.6), "私聊功能",
     "• 独立会话窗口，实时收发\n• 离线投递：对方离线消息仍入库，上线后可查\n• 未读红点：打开窗口自动标记已读\n• 会话历史独立保存，支持关键字搜索\n• 任意一方清空记录不影响对方（双份行模型）", body_size=11.5)
card(s, Inches(9.0), Inches(4.1), Inches(3.85), Inches(2.45), "管理员功能",
     "• 踢出在线用户、封禁账号（删除账号及私聊数据）\n• 清空公共聊天记录（向全体广播）\n• 查看任意两名用户的私聊记录\n• 撤回任意用户的公共消息", body_size=11.5)
footer(s, 10)

# ---------- 11. 系统实现：技术栈 ----------
s = add_slide()
page_header(s, "系统实现 —— 技术栈与工程化")
card(s, Inches(0.6), Inches(1.3), Inches(5.9), Inches(2.5), "开发环境",
     "• 语言：Java（JDK 17+）\n• 界面：Swing / AWT（Graphics2D 自绘）\n• 网络：java.net Socket / ServerSocket\n• 数据库：MySQL 8.0 + JDBC\n• 图片处理：ImageIO / Base64")
card(s, Inches(6.8), Inches(1.3), Inches(5.9), Inches(2.5), "工程化实践",
     "• 13 个 Java 类，约 3.4 万行实现（含 UI）\n• IDEA + Git 团队协作\n• build.bat 一键编译打包，发布 release 产物（LanChat-1.1.0）\n• config.properties 配置数据库，jar 同目录优先查找\n• 密码等敏感配置已被 .gitignore 排除")
card(s, Inches(0.6), Inches(4.0), Inches(5.9), Inches(2.6), "安全设计",
     "• 换行注入双重过滤（发送侧 + 服务端 sanitize）\n• 消息限长 2000 字符，图片整体校验\n• 图片 / 头像魔数校验（JPEG / PNG 文件头）\n• 管理命令统一 isAdmin 前置校验\n• 未登录连接忽略一切命令（防 DDoS）\n• putIfAbsent 原子注册防重复登录")
card(s, Inches(6.8), Inches(4.0), Inches(5.9), Inches(2.6), "健壮性设计",
     "• 断开清理统一收敛 disconnect()，幂等防重复广播\n• 协议解析 splitFixed + parseLong 兜底容错\n• 历史回放异常分支也发送 END，客户端不卡加载\n• 自聊拦截、超长消息截断、空白消息丢弃")
footer(s, 11)

# ---------- 12. 系统测试：环境与统计 ----------
s = add_slide()
page_header(s, "系统测试 —— 测试环境与执行情况")
card(s, Inches(0.6), Inches(1.25), Inches(5.9), Inches(2.35), "测试环境",
     "• 操作系统：Windows 11\n• 运行时：JDK 17（OpenJDK 17）\n• 数据库：MySQL 8.0，localhost:3306/lanchat\n• 服务端：chatServer 监听 8080\n• 客户端：Swing 桌面程序，双实例并发联测")
card(s, Inches(6.8), Inches(1.25), Inches(5.9), Inches(2.35), "测试策略",
     "• 客户端功能测试：启动服务端与客户端，按功能域逐项操作界面验证\n• 覆盖六大功能域 34 个用例、37 项验证点\n• 含 5 个异常与安全用例：特殊字符消息、非法图片文件、越权操作、伪造消息 ID、未登录访问\n• 独立测试账号 + 双客户端交互验证 + 测试后数据核查")
# 统计条
stats = [
    ("34", "测试用例"), ("37", "验证点"), ("100%", "通过率"), ("0", "遗留缺陷"),
]
for i, (v, t) in enumerate(stats):
    bx = Inches(0.6 + i * 3.15)
    rect(s, bx, Inches(3.85), Inches(2.85), Inches(1.5), LIGHT,
         line=RGBColor(0xD5, 0xE2, 0xF0), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, bx, Inches(4.02), Inches(2.85), Inches(0.7), v, size=32, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, bx, Inches(4.72), Inches(2.85), Inches(0.4), t, size=13, color=GRAY, align=PP_ALIGN.CENTER)
# 分域统计表
tbl = s.shapes.add_table(7, 3, Inches(0.6), Inches(5.55), Inches(12.2), Inches(1.4)).table
tbl.columns[0].width = Inches(5.2)
tbl.columns[1].width = Inches(3.0)
tbl.columns[2].width = Inches(4.0)
domain = [
    ("认证与登录", "6", "100%"), ("公共消息", "5", "100%"), ("私聊", "6", "100%"),
    ("消息撤回与 @ 提醒", "4", "100%"), ("管理员功能与权限", "6", "100%"),
    ("图片、头像、搜索与健壮性", "10", "100%"),
]
for j, h in enumerate(["功能域", "用例数", "通过率"]):
    c = tbl.cell(0, j)
    c.fill.solid(); c.fill.fore_color.rgb = NAVY
    c.text_frame.text = h
    for p in c.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
for i, (d, n, rate) in enumerate(domain):
    for j, v in enumerate([d, n, rate]):
        c = tbl.cell(i + 1, j)
        c.fill.solid()
        c.fill.fore_color.rgb = LIGHT if i % 2 == 0 else WHITE
        c.text_frame.text = v
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(11.5); r.font.color.rgb = DARK; r.font.name = FONT
footer(s, 12)

# ---------- 13. 系统测试：结论 ----------
s = add_slide()
page_header(s, "系统测试 —— 测试结论")
concl = [
    ("功能正确性", "认证、公共消息、私聊（含离线投递）、撤回、@提醒、管理功能、图片头像、搜索全部符合需求约定；协议字段解析（内容含冒号、时间戳数字格式）与设计一致。"),
    ("安全性", "换行注入、伪造图片、越权命令、伪造消息 ID、未登录注入等攻击场景均被服务端防线拦截：sanitize 清洗、魔数校验、isAdmin 权限校验、putIfAbsent 原子注册等机制有效。"),
    ("健壮性", "超长消息截断、空白消息丢弃、未知命令容错、异常断连清理正常；多客户端并发下服务端运行稳定，无内存泄漏迹象。"),
    ("接口一致性", "网络层与数据层、界面层衔接正常：消息 ID、未读数、双份行模型等跨模块约定经实测与设计一致。"),
]
for i, (t, d) in enumerate(concl):
    by = Inches(1.3 + i * 1.35)
    rect(s, Inches(0.6), by, Inches(12.15), Inches(1.12), LIGHT, line=RGBColor(0xD5, 0xE2, 0xF0),
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, Inches(0.6), by, Inches(1.75), Inches(1.12), NAVY)
    txt(s, Inches(0.6), by + Inches(0.3), Inches(1.75), Inches(0.5), t,
        size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(2.55), by + Inches(0.1), Inches(10.0), Inches(0.95), d,
        size=11.8, color=DARK, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12)
txt(s, Inches(0.6), Inches(6.75), Inches(12.2), Inches(0.4),
    "结论：网络通信与服务端模块功能完备、安全机制有效、运行稳定，通过全部测试，可以交付。",
    size=13.5, color=GREEN, bold=True)
footer(s, 13)

# ---------- 14. 结束页 ----------
s = add_slide()
rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
rect(s, 0, Inches(3.3), SLIDE_W, Inches(0.03), ORANGE)
txt(s, 0, Inches(2.6), SLIDE_W, Inches(0.8), "感谢观看，请批评指正",
    size=34, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, 0, Inches(3.65), SLIDE_W, Inches(0.45), "局域网聊天室 · 小组成员：赵翔 · 胡鸣 · 茹桂堂",
    size=15, color=RGBColor(0xBD, 0xD7, 0xEE), align=PP_ALIGN.CENTER)
txt(s, 0, Inches(4.25), SLIDE_W, Inches(0.4), "答辩流程：一名同学主汇报，其他两名同学补充，并记录老师提问",
    size=12, color=RGBColor(0x9D, 0xC3, 0xE6), align=PP_ALIGN.CENTER)

OUT_PPTX = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "汇报答辩PPT-即时通讯工具.pptx")
prs.save(OUT_PPTX)
print("PPT 已生成:", os.path.abspath(OUT_PPTX), "页数:", len(prs.slides._sldIdLst))
